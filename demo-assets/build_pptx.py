#!/usr/bin/env python3
"""Generate an editable SQD-APP demonstration deck (PPTX) from the same
content as SQD-APP-Demo.html. Text is authored as native PowerPoint text
boxes/tables/shapes (fully editable); screenshots and flowcharts are embedded
as pictures. Palette follows DESIGN.md."""
import struct, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
SCR = os.path.join(HERE, "screens")

# ---- palette ----
BLUE=RGBColor(0x25,0x63,0xEB); BLUE_S=RGBColor(0xEF,0xF6,0xFF)
AMBER=RGBColor(0xD9,0x77,0x06); AMBER_S=RGBColor(0xFF,0xFB,0xEB)
RED=RGBColor(0xDC,0x26,0x26); RED_S=RGBColor(0xFE,0xF2,0xF2)
EMER=RGBColor(0x05,0x96,0x69); EMER_S=RGBColor(0xEC,0xFD,0xF5)
INK=RGBColor(0x1E,0x29,0x3B); INK2=RGBColor(0x47,0x55,0x69); MUT=RGBColor(0x94,0xA3,0xB8)
BASE=RGBColor(0xF8,0xFA,0xFC); CARD=RGBColor(0xFF,0xFF,0xFF); BORD=RGBColor(0xE2,0xE8,0xF0)
DARK=RGBColor(0x0B,0x12,0x20); WHITE=RGBColor(0xFF,0xFF,0xFF)
SANS="Geist Sans"; MONO="Geist Mono"

EMU_IN=914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def png_size(path):
    with open(path,'rb') as f:
        head=f.read(24)
    w,h=struct.unpack('>II',head[16:24])
    return w,h

def slide(bg=BASE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,x,y,w,h):
    return s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))

def setp(p,text,size,color,bold=False,italic=False,font=SANS,align=PP_ALIGN.LEFT,spacing=1.0):
    p.alignment=align
    if spacing: p.line_spacing=spacing
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.color.rgb=color; r.font.bold=bold
    r.font.italic=italic; r.font.name=font
    return p

def add_text(s,x,y,w,h,lines):
    """lines: list of dicts {t,size,color,bold,italic,font,align,space_before}"""
    tb=box(s,x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,ln in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        setp(p,ln.get('t',''),ln.get('size',14),ln.get('color',INK2),
             ln.get('bold',False),ln.get('italic',False),ln.get('font',SANS),
             ln.get('align',PP_ALIGN.LEFT),ln.get('space',1.15))
        if ln.get('space_before'): p.space_before=Pt(ln['space_before'])
        if ln.get('bullet'): _bullet(p)
    return tb

def _bullet(p):
    pPr=p._pPr if p._pPr is not None else p.get_or_add_pPr()
    buChar=pPr.makeelement(qn('a:buChar'),{'char':'•'});
    buFont=pPr.makeelement(qn('a:buFont'),{'typeface':'Arial'})
    pPr.append(buFont); pPr.append(buChar)
    pPr.set('indent','-137160'); pPr.set('marL','137160')

def rect(s,x,y,w,h,fill=CARD,line=BORD,line_w=1.0,rounded=True):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(line_w)
    shp.shadow.inherit=False
    try: shp.adjustments[0]=0.06
    except Exception: pass
    return shp

def kicker_title(s,kicker,title,sub=None,color=BLUE,tcolor=INK):
    add_text(s,0.7,0.42,12,0.4,[{'t':kicker.upper(),'size':12,'color':color,'bold':True}])
    add_text(s,0.7,0.72,12,1.0,[{'t':title,'size':30,'color':tcolor,'bold':True,'space':1.05}])
    if sub:
        add_text(s,0.7,1.55,11.9,0.7,[{'t':sub,'size':14,'color':INK2,'space':1.2}])

def footer(s,n):
    add_text(s,0.7,7.05,8,0.3,[{'t':'SQD-APP · Aviation MRO Quality Assurance','size':10,'color':MUT,'bold':True}])
    add_text(s,12.0,7.05,1.0,0.3,[{'t':str(n),'size':10,'color':MUT,'align':PP_ALIGN.RIGHT}])

def picture_contain(s,path,x,y,w,h,border=True,shadow=True):
    iw,ih=png_size(path); ar=iw/ih; boxar=w/h
    if ar>boxar: dw=w; dh=w/ar
    else: dh=h; dw=h*ar
    px=x+(w-dw)/2; py=y+(h-dh)/2
    pic=s.shapes.add_picture(path,Inches(px),Inches(py),Inches(dw),Inches(dh))
    if border:
        pic.line.color.rgb=BORD; pic.line.width=Pt(1)
    return pic

def chip(s,x,y,text,fg,bg,w=1.4):
    c=rect(s,x,y,w,0.32,fill=bg,line=None);
    tf=c.text_frame; tf.word_wrap=False; tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    p=tf.paragraphs[0]; setp(p,text,10,fg,bold=True,align=PP_ALIGN.CENTER)
    return c

def card_with_bullets(s,x,y,w,h,head,bullets,head_color=INK):
    rect(s,x,y,w,h,fill=CARD,line=BORD)
    add_text(s,x+0.25,y+0.18,w-0.5,0.4,[{'t':head,'size':15,'color':head_color,'bold':True}])
    lines=[{'t':b,'size':12.5,'color':INK2,'bullet':True,'space':1.12,'space_before':4} for b in bullets]
    add_text(s,x+0.25,y+0.62,w-0.5,h-0.8,lines)

n=0
def N():
    global n; n+=1; return n

# ---------- 1 TITLE ----------
s=slide(DARK); N()
add_text(s,0.9,2.4,11,0.4,[{'t':'AVIATION MRO · QUALITY ASSURANCE & CONTROL','size':13,'color':RGBColor(0x60,0xA5,0xFA),'bold':True}])
add_text(s,0.9,2.8,11,1.1,[{'t':'SQD-APP','size':54,'color':WHITE,'bold':True}])
add_text(s,0.9,3.95,10.5,1.2,[{'t':'A full-stack platform for dynamic audit templates, work packages, task execution, findings and corrective action — built to enforce compliance and keep an immutable, auditable record of every quality event.','size':17,'color':RGBColor(0x94,0xA3,0xB8),'space':1.25}])
add_text(s,0.9,5.4,11,0.4,[{'t':'Audience: Staff · Managers · Directors      Aesthetic: “The Technical Manual”      Live demo data: 30 WPs · 150 Tasks · 27 Findings','size':13,'color':RGBColor(0xCB,0xD5,0xE1)}])

# ---------- 2 PURPOSE ----------
s=slide(); kicker_title(s,'Introduction','What it is, and the design philosophy',
  'An internal tool for a high-stakes, compliance-driven environment where errors have safety consequences.')
card_with_bullets(s,0.7,2.35,5.85,4.2,'🎯  Product purpose',[
  'Manage audit templates, task assignment, inspections, findings and work packages across the MRO org.',
  'Enforce compliance workflows and surface issues before they become safety incidents.',
  'Success = inspectors work without friction; managers trust the data they see.'])
card_with_bullets(s,6.75,2.35,5.85,4.2,'🧭  “The Technical Manual”',[
  'Authority through restraint — structure signals safety.',
  'Zero ambiguity on status — readable in under one second.',
  'One Signal Rule — Blue means “act here”; status colours fire only on live conditions.',
  'A screen free of amber / red is itself a signal: the system is healthy.'])
footer(s,n)

# ---------- 3 TECH STACK ----------
s=slide(); kicker_title(s,'Introduction','Technology & platform')
card_with_bullets(s,0.7,2.15,3.9,4.4,'Frontend',[
  'Next.js 16 (App Router, React 19)','Tailwind CSS v4, Zustand auth','TypeScript 5 strict, Axios'])
card_with_bullets(s,4.72,2.15,3.9,4.4,'Backend',[
  'Express 5, Node, TypeScript strict','Prisma v6 ORM + Postgres adapter','JWT auth, bcrypt, single-session','Jest + Supertest (~500 tests)'])
card_with_bullets(s,8.74,2.15,3.9,4.4,'Data & compliance',[
  'PostgreSQL — status machines enforced by DB CHECK constraints','Soft-delete on all compliance records','Dual audit: AuditLog + operational feed'])
footer(s,n)

# ---------- 4 ERD ----------
s=slide(); kicker_title(s,'Architecture','Data model — two clusters around the Task',
  'Not a linear pipeline: Task is the hub, and Finding is both born from a task and resolved by new tasks.')
def erd_node(s,x,y,w,t,k,fill=CARD,line=BORD,tcol=INK):
    rect(s,x,y,w,0.72,fill=fill,line=line,line_w=1.4)
    add_text(s,x,y+0.10,w,0.35,[{'t':t,'size':14,'color':tcol,'bold':True,'align':PP_ALIGN.CENTER}])
    add_text(s,x,y+0.42,w,0.25,[{'t':k,'size':10,'color':MUT,'font':MONO,'align':PP_ALIGN.CENTER}])
add_text(s,0.9,2.5,4,0.3,[{'t':'TASK HUB','size':11,'color':MUT,'bold':True}])
erd_node(s,1.6,2.9,3.0,'Template','formSchema · draftSchema')
add_text(s,1.6,3.66,3.0,0.25,[{'t':'│ 1 : N  (schema snapshot at creation)','size':10,'color':INK2,'font':MONO,'align':PP_ALIGN.CENTER}])
erd_node(s,1.6,3.95,3.0,'Task','schemaSnapshot · status(9)',fill=BLUE_S,line=BLUE,tcol=BLUE)
add_text(s,1.6,4.71,3.0,0.25,[{'t':'▲ N : 1   wpId (optional)','size':10,'color':INK2,'font':MONO,'align':PP_ALIGN.CENTER}])
erd_node(s,1.6,5.0,3.0,'WorkPackage','autoGenTemplateId → Template')
add_text(s,7.4,2.5,4,0.3,[{'t':'FINDING CLUSTER','size':11,'color':MUT,'bold':True}])
erd_node(s,8.1,2.9,3.0,'Task','source & follow-up',fill=BLUE_S,line=BLUE,tcol=BLUE)
add_text(s,8.1,3.66,3.0,0.25,[{'t':'│ sourceTaskId ▼   ▲ parentFindingId','size':10,'color':INK2,'font':MONO,'align':PP_ALIGN.CENTER}])
erd_node(s,8.1,3.95,3.0,'Finding','severity · status(5)',fill=RED_S,line=RED,tcol=RED)
add_text(s,8.1,4.72,3.0,0.25,[{'t':'└─ RCA · CAPA · ATA · Hazard tags','size':10,'color':INK2,'font':MONO,'align':PP_ALIGN.CENTER}])
erd_node(s,8.1,5.05,1.45,'RcaInvestigation','')
erd_node(s,9.65,5.05,1.45,'CapaAction','')
add_text(s,0.7,6.55,12,0.4,[{'t':'Every arrow is a real foreign key in schema.prisma · Task.templateId, Task.wpId, Finding.sourceTaskId, Task.parentFindingId.','size':11,'color':MUT,'italic':True}])
footer(s,n)

# ---------- 5 SCHEMA ----------
s=slide(); kicker_title(s,'Architecture','Schema & the three immutability guarantees')
card_with_bullets(s,0.7,2.1,3.9,1.9,'📸  Schema snapshot',['Each Task stores an immutable schemaSnapshot at creation — template edits never alter an in-flight task.'])
card_with_bullets(s,4.72,2.1,3.9,1.9,'🗑️  Soft delete',['Records with deletedAt are never physically removed. Evidence persists — a compliance requirement.'])
card_with_bullets(s,8.74,2.1,3.9,1.9,'📝  Dual audit',['Every state change writes to both the system-wide AuditLog and the operational feed.'])
add_text(s,0.7,4.2,6,0.3,[{'t':'CORE MODELS (30+ TOTAL)','size':11,'color':MUT,'bold':True}])
models=[('Template','Audit form + branching'),('Task','9-status execution unit'),('WorkPackage','Timeframed container'),('Finding','Defect / quality event'),
('RcaInvestigation','5-Whys / MEDA'),('CapaAction','Corrective / preventive'),('TemplateSet','Reusable bundle'),('WpBlueprint','Recurring WP recipe'),
('FeedPost','Unified 5-scope feed'),('EscalationFlag','Escalation lifecycle'),('TimeEntry','Budget vs actual, immutable'),('AuditLog','Immutable compliance log')]
gx,gy=0.7,4.55; cw,ch=3.02,0.62
for i,(nm,de) in enumerate(models):
    r=i//4; c=i%4; x=gx+c*(cw+0.06); y=gy+r*(ch+0.08)
    rect(s,x,y,cw,ch,fill=CARD,line=BORD)
    add_text(s,x+0.14,y+0.06,cw-0.2,0.3,[{'t':nm,'size':12.5,'color':INK,'bold':True,'font':MONO}])
    add_text(s,x+0.14,y+0.32,cw-0.2,0.25,[{'t':de,'size':10,'color':MUT}])
footer(s,n)

# ---------- 6 RBAC TABLE ----------
s=slide(); kicker_title(s,'Architecture','RBAC — view-transparent, action-scoped',
  'Everyone can SEE everything; only ACTING is restricted by role, privilege and division.')
rows=[('Capability','Dir','Adm','Mgr','Staff*'),
('Create tasks','✅','✅','✅','—'),
('Assign across any division','✅','✅','—','—'),
('Assign within own division','✅','✅','✅','—'),
('Review tasks','✅ any','—','✅ div','—'),
('Close / reopen Work Package','✅','✅','creator','—'),
('Assign users to WPs','✅','✅','✅','—'),
('Review findings (severity/close)','✅','—','✅ div','—'),
('Action escalations','✅','✅','✅ scoped','—'),
('Analytics','✅','✅','✅ div','—'),
('Manage privilege config','—','✅','—','—')]
tw,th=9.3,4.35; tx,ty=0.7,2.35
gtbl=s.shapes.add_table(len(rows),5,Inches(tx),Inches(ty),Inches(tw),Inches(th)).table
gtbl.columns[0].width=Inches(4.5)
for c in range(1,5): gtbl.columns[c].width=Inches((tw-4.5)/4)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell=gtbl.cell(ri,ci); cell.margin_top=Pt(2); cell.margin_bottom=Pt(2)
        cell.margin_left=Pt(6); cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=cell.text_frame.paragraphs[0]; p.text=val
        r=p.runs[0]; r.font.name=SANS; r.font.size=Pt(11.5 if ri else 12)
        if ri==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=INK
            r.font.color.rgb=WHITE; r.font.bold=True
            if ci>0: p.alignment=PP_ALIGN.CENTER
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb=CARD if ri%2 else BASE
            r.font.color.rgb=INK if ci==0 else INK2
            r.font.bold=(ci==0)
            if ci>0: p.alignment=PP_ALIGN.CENTER
add_text(s,10.2,2.4,2.6,4.4,[
  {'t':'* Staff / Group Leader / Senior Advisor hold no privilege keys by default — they act via hardcoded relationship grants (self-assign, perform, raise findings, book time).','size':11,'color':INK2,'space':1.2,'space_before':0},
  {'t':'Privileges are DB-driven (PrivilegeConfig); division-scope, relationship bypasses and the Director-approval gate stay hardcoded.','size':11,'color':INK2,'space':1.2,'space_before':10},
  {'t':'Full detail: RBAC_RULES.md','size':11,'color':BLUE,'bold':True,'space_before':10}])
footer(s,n)

# ---------- 7 LIFECYCLE OVERVIEW ----------
s=slide(); kicker_title(s,'The Centrepiece','The end-to-end lifecycle')
steps=[('01','Template','Branching audit form'),('02','Work Package','Auto-generate tasks'),('03','Execution','Staff perform; raise findings'),
('04','Findings','Severity, follow-ups, RCA/CAPA'),('05','Time + Review','Budget vs actual; approve/reject'),('06','Closure','Director closes; analytics')]
sx=0.7; sw=1.98;
for i,(nm,t,d) in enumerate(steps):
    x=sx+i*(sw+0.06)
    rect(s,x,2.15,sw,1.15,fill=CARD,line=BORD)
    add_text(s,x+0.12,2.24,sw-0.2,0.25,[{'t':nm,'size':11,'color':BLUE,'bold':True,'font':MONO}])
    add_text(s,x+0.12,2.5,sw-0.2,0.3,[{'t':t,'size':13,'color':INK,'bold':True}])
    add_text(s,x+0.12,2.82,sw-0.2,0.4,[{'t':d,'size':10.5,'color':INK2,'space':1.05}])
picture_contain(s,os.path.join(SCR,'director-01-dashboard.png'),0.7,3.5,11.9,3.2)
add_text(s,0.7,6.72,11.9,0.3,[{'t':'Director “Operations Overview” — live seeded data: 53 pending tasks (19 overdue), 16 findings, active work packages, live feed.','size':11,'color':MUT,'italic':True,'align':PP_ALIGN.CENTER}])
footer(s,n)

# ---------- screenshot slides helper ----------
def shot_slide(kicker,title,bullets,img,cap,two=None):
    s=slide(); kicker_title(s,kicker,title)
    lines=[{'t':b,'size':12.5,'color':INK2,'bullet':True,'space':1.12,'space_before':4} for b in bullets]
    add_text(s,0.7,1.75,11.9,1.0,lines)
    top=2.9
    if two:
        picture_contain(s,os.path.join(SCR,img),0.7,top,5.9,3.6)
        picture_contain(s,os.path.join(SCR,two),6.75,top,5.9,3.6)
    else:
        picture_contain(s,os.path.join(SCR,img),0.7,top,11.9,3.6)
    add_text(s,0.7,6.6,11.9,0.35,[{'t':cap,'size':11,'color':MUT,'italic':True,'align':PP_ALIGN.CENTER}])
    footer(s,n); return s

# 8 Manager WP autogen
N(); shot_slide('Step 2A · Manager','Work Package setup & automatic task generation',
 ['Any WP type can auto-generate tasks — from a single Template, a saved Template Set, or an inline set.',
  'Two modes: Single-shot (spawn once) and Repeat (every N days). First batch fires the moment the WP is saved.'],
 'manager-02-wp-new-autogen.png','“Automatic Task Generation — spawn tasks from a template automatically, once or on a repeating cadence.”')
# 9 Sets & Blueprints
N(); shot_slide('Step 2A · Manager','Reusable recipes — Template Sets & WP Blueprints',
 ['Standardise once, launch repeatedly, or let a nightly cron auto-launch recurring routine WPs.'],
 'manager-03-template-sets.png','Template Sets (ordered bundles) · WP Blueprints (pre-filled WP + auto-gen + recurrence).','manager-04-wp-blueprints.png')
# 10 Staff tasks + unassigned
N(); shot_slide('Step 2B · Staff','Task execution & the self-serve pool',
 ['Tasks arrive assigned, or sit Unassigned in a pool eligible Staff self-serve via “Perform This Task”.',
  'Role-adaptive sidebar: Staff see a focused nav (no Analytics / Template Builder / Escalations).'],
 'staff-02-tasks.png','Task list with the status-badge column · an unassigned task ready to be picked up.','staff-03-task-unassigned.png')
# 11 Staff execution form
N(); shot_slide('Step 2B · Staff','Executing a task — snapshot form & status vocabulary',
 ['The form renders from the task’s immutable snapshot — it can never drift from a later template edit.',
  'Templates support Google-Forms-style branching. A Raise Finding action appears when the template allows it.'],
 'staff-04-task-inprogress-form.png','Task detail — template link, issuer/assignee, WP, deadline, est. hours, progress editor, activity feed.')
# 12 Findings
N(); shot_slide('Step 2C · Staff → Manager','Raising & routing findings',
 ['Any user can raise a finding (if the template allows and the task isn’t final).',
  'A Manager/Director assigns severity — Observation / Level 1 / Level 2 — and a due date; visibility follows RBAC.'],
 'director-04-findings.png','Findings register — status & severity badges read in under one second.')
# 13 RCA/CAPA
N(); shot_slide('Differentiator · Headline','Root Cause & Corrective Action (RCA / CAPA)',
 ['RCA (1:1): 5-Whys, MEDA contributing factors, or OTHER — concluding in a cause code. CAPA: corrective & preventive.',
  'Taxonomy & traceability: ATA chapters, hazard tags, finding links, response actions. Trend detection fires a recurrence banner.'],
 'director-05-finding-hero-rca-capa.png','Hero finding FND-000101 — “Recurrent pattern detected” banner, ATA chapter, hazard tags, Level 2, full lifecycle.')
# 14 Escalation
N(); shot_slide('Differentiator · Headline','Unified Feed & the Escalation loop',
 ['One feed, five scopes (Task · WP · Division · Org · Finding). Reading is open; any user can flag a comment upward.',
  'Escalation posts a card at the target + info cards at every level in between. Only Directors/Admins & scoped Managers can action.'],
 'director-06-escalations.png','The Escalations queue — full history (Pending / Actioned / Dismissed), reusing existing workflows.')

# 15 Time & review
s=slide(); N(); kicker_title(s,'Step 2D · Staff → Manager','Time booking & review')
card_with_bullets(s,0.7,2.15,3.9,4.3,'⏱️  Budget vs actual',[
  'At a final state, the assignee logs actual hours (+ collaborators).',
  'If the template had estimatedHours, a budget-vs-actual badge appears.',
  'Exceeding 120% forces an over-budget reason.',
  'Each submission writes an append-only TimeEntry — immutable trail.'])
card_with_bullets(s,4.72,2.15,3.9,4.3,'✔️  Review rights (precision)',[
  'Approve / Reject / Follow-up = Issuer + Director + Managers of the same Division.',
  'The assignee can NEVER review their own work (segregation of duties).',
  'QN tasks require Director approval.',
  'Reassign allowed at any non-final stage (reason; data preserved).'])
card_with_bullets(s,8.74,2.15,3.9,4.3,'📊  Feeds Analytics',[
  'Managers & Directors see time-efficiency trends and per-staff performance.',
  'Scoped to division or system-wide.'])
footer(s,n)

# 16 Analytics
N(); shot_slide('Step 2E · Director','Oversight & analytics',
 ['Directors close tasks & work packages; global action rights across all divisions.',
  'Analytics: Time Efficiency, Findings and Personnel tabs — avg rating, tasks rated, efficiency multiplier.'],
 'director-07-analytics.png','Staff Performance — average rating, tasks rated, and efficiency (green = ahead of estimate, red = behind).')

# 17 Per-role summary
s=slide(); N(); kicker_title(s,'Summary','What each role does')
def role_card(x,head,sub,items):
    rect(s,x,2.15,3.9,4.5,fill=CARD,line=BORD)
    add_text(s,x+0.25,2.32,3.5,0.35,[{'t':head,'size':17,'color':INK,'bold':True}])
    add_text(s,x+0.25,2.72,3.5,0.3,[{'t':sub,'size':11,'color':MUT}])
    add_text(s,x+0.25,3.1,3.5,3.4,[{'t':it,'size':12,'color':INK2,'bullet':True,'space':1.12,'space_before':4} for it in items])
role_card(0.7,'Staff','Perform the work',['Pick up unassigned tasks','Execute snapshot forms','Raise findings from tasks','Log actual hours','Comment & flag on feeds','View all work (transparency)'])
role_card(4.72,'Manager','Plan & review (division)',['Create WPs, auto-generate tasks','Template Sets & Blueprints','Assign within division','Approve / reject / follow-up','Triage findings + follow-ups','Action escalations in scope','Division analytics'])
role_card(8.74,'Director / Admin','Global oversight',['Global assign & review','Close tasks & work packages','Action any escalation; moderate','System-wide RCA/CAPA sign-off','System-wide analytics','Configure privileges & settings'])
footer(s,n)

# 18 Swimlane
s=slide(); N(); kicker_title(s,'Workflow','Swimlane — task lifecycle across roles (all cases)')
picture_contain(s,os.path.join(HERE,'workflow-swimlane.png'),0.7,1.7,11.9,5.0,border=False)
add_text(s,0.7,6.72,11.9,0.3,[{'t':'System · Staff · Manager · Director lanes — self-assign, review (segregation of duties), finding loop, reassign, closure. Full-res PNG in demo-assets/.','size':10.5,'color':MUT,'italic':True,'align':PP_ALIGN.CENTER}])
footer(s,n)

# 19 Escalation flow
s=slide(); N(); kicker_title(s,'Workflow','Swimlane — Unified Feed & Escalation loop')
picture_contain(s,os.path.join(HERE,'workflow-escalation.png'),0.7,1.7,11.9,5.0,border=False)
add_text(s,0.7,6.72,11.9,0.3,[{'t':'Any user flags → system posts cards → Director/Admin or scoped Manager actions → connects back to the Findings / Task loops.','size':10.5,'color':MUT,'italic':True,'align':PP_ALIGN.CENTER}])
footer(s,n)

# 20 RBAC key rules
s=slide(); N(); kicker_title(s,'Governance','RBAC — the rules that matter most')
card_with_bullets(s,0.7,2.1,5.85,4.5,'Assignment & division scope',[
  'Assign / reassign only to your OWN division — unless you hold assign_any (Director/Admin).',
  'Self-assign only from your own division’s unassigned pool.',
  'WP membership lets you create/assign within that WP — still own-division-locked.',
  'Multi-division WP: only a Director/Admin can add out-of-division members; routing work still stays division-scoped per actor.'])
card_with_bullets(s,6.75,2.1,5.85,4.5,'Review, findings & escalation',[
  'Review = Issuer + Director + Manager (same division). Admin is NOT a reviewer.',
  'The person who performed a task can never review it.',
  'QN tasks require Director approval (overrides Issuer).',
  'Manager reviews findings only in their division scope.',
  'Escalations: Director/Admin any; Manager own-division WP/Div + all Org.'])
add_text(s,0.7,6.75,12,0.3,[{'t':'Complete, code-verified list: RBAC_RULES.md','size':11,'color':BLUE,'bold':True}])
footer(s,n)

# 21 Live-demo appendix
s=slide(); N(); kicker_title(s,'Appendix','Live-demo runbook — exact click paths')
steps=[
 ('1 · Groundwork','Open the deck slides 4–6 (ERD, schema, RBAC). Frame “view-transparent, action-scoped”.'),
 ('2A · Manager','Login VAE00061. Work Packages → New → toggle Automatic Task Generation → pick a Template → Create. Open the WP: tasks already spawned.'),
 ('2B · Staff','Login VAE00051. Tasks → Unassigned tab → “Perform This Task”. Open the task → fill the branching form → note status badge → Submit.'),
 ('2C · Findings','On a task with findings enabled → Raise Finding. As Manager: open Findings → set severity + due → Generate follow-up task.'),
 ('3 · RCA/CAPA','Findings → open FND-000101 (hero) → show recurrence banner, ATA/hazard tags, RCA (5-Whys/MEDA) + CAPA.'),
 ('3 · Escalation','Any feed → flag a comment → header bell → Escalations page → action it (Acknowledge / Create Task / Disseminate).'),
 ('2D · Time+Review','On a Closed task → book hours (budget vs actual). As Issuer/Manager → Review → Approve / Reject / Follow-up.'),
 ('2E · Director','Login VAE00071. Close the Work Package → Analytics (Time Efficiency / Personnel).')]
y=2.05
for h,d in steps:
    rect(s,0.7,y,11.9,0.56,fill=CARD,line=BORD)
    add_text(s,0.85,y+0.08,2.5,0.4,[{'t':h,'size':12,'color':BLUE,'bold':True}])
    add_text(s,3.35,y+0.08,9.1,0.45,[{'t':d,'size':11.5,'color':INK2,'space':1.05}])
    y+=0.62
add_text(s,0.7,y+0.02,12,0.3,[{'t':'Accounts (password Demo@12345): Director VAE00071 · Manager VAE00061 · Staff VAE00051.  Set DISABLE_RATE_LIMIT=true for rapid multi-role logins.','size':10.5,'color':MUT,'italic':True}])
footer(s,n)

# 22 Close
s=slide(DARK); N()
add_text(s,0.9,2.5,11.5,0.4,[{'t':'WRAP-UP','size':13,'color':RGBColor(0x60,0xA5,0xFA),'bold':True}])
add_text(s,0.9,2.9,11.5,1.0,[{'t':'Everything traces back to the audit log','size':36,'color':WHITE,'bold':True}])
add_text(s,0.9,4.0,11,1.0,[{'t':'Every status change, time entry, escalation and closure shown today landed in the immutable AuditLog — the compliance backbone that lets managers trust the data.','size':16,'color':RGBColor(0x94,0xA3,0xB8),'space':1.25}])
add_text(s,0.9,5.6,11,0.4,[{'t':'9 task statuses (DB-enforced)   ·   5 feed scopes (one model)   ·   6 roles (DB-driven privileges)   ·   30+ schema models','size':14,'color':RGBColor(0xCB,0xD5,0xE1)}])

out=os.path.join(HERE,'SQD-APP-Demo.pptx')
prs.save(out)
print('Saved',out,'-',len(prs.slides.__iter__.__self__._sldIdLst),'slides')
